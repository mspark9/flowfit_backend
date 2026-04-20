"""
경쟁사 리서치 PPTX 생성 서비스
슬라이드 구성:
  1   표지
  2~N 회사별 동향
  N+1 종합 시사점
  N+2 시가총액 비교 차트     ┐ financial_data가 있을 때만
  N+3 매출·영업이익 비교 차트 ┤ 각 차트는 데이터 0개면 스킵
  N+4 PER·PBR 비교 차트     ┘
"""
import io
from datetime import datetime


def generate_research_pptx(data: dict) -> bytes:
    """
    data 구조:
      companies:     [{ name, analysis: { category: text }, articles: [...] }]
      summary:       str
      categories:    [str]
      financial_data: { ticker: { metric: value } }   # 선택
      ticker_map:    { company_name: { ticker, exchange, found } }  # 선택
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # ── 색상 팔레트 ────────────────────────────────────────────
    AMBER      = RGBColor(0xD9, 0x77, 0x06)
    AMBER_DARK = RGBColor(0x92, 0x40, 0x09)
    AMBER_LITE = RGBColor(0xFF, 0xF7, 0xED)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    DARK       = RGBColor(0x1F, 0x29, 0x37)
    GRAY       = RGBColor(0x6B, 0x72, 0x80)
    GRAY_LITE  = RGBColor(0xF3, 0xF4, 0xF6)

    CHART_C1   = RGBColor(0xD9, 0x77, 0x06)  # amber-600
    CHART_C2   = RGBColor(0xF5, 0x9E, 0x0B)  # amber-400

    companies  = data.get("companies", [])
    summary    = data.get("summary", "")
    categories = data.get("categories", [])
    today      = datetime.now().strftime("%Y년 %m월 %d일")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    # ── 헬퍼: 텍스트박스 ──────────────────────────────────────
    def add_textbox(slide, left, top, width, height, text,
                    font_size=14, bold=False, color=DARK,
                    align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf    = txBox.text_frame
        tf.word_wrap = wrap
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = color
        return txBox

    def add_rect(slide, left, top, width, height, fill_color, line_color=None):
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    # ──────────────────────────────────────────────────────────
    # 슬라이드 1: 표지
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, int(H * 0.65), AMBER)
    add_rect(slide, 0, int(H * 0.65), W, int(H * 0.35), WHITE)

    add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
                "Front-Office  |  전략/기획팀",
                font_size=12, color=RGBColor(0xFF, 0xE0, 0x9A), align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                "경쟁사 동향 리서치 보고서",
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    company_names = " · ".join(c["name"] for c in companies)
    add_textbox(slide, Inches(1), Inches(3.1), Inches(11), Inches(0.6),
                company_names,
                font_size=16, color=RGBColor(0xFF, 0xEC, 0xB3), align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
                today, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)
    cat_text = "  |  ".join(categories)
    add_textbox(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.4),
                f"조사 카테고리: {cat_text}", font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────
    # 슬라이드 2~N: 회사별 동향
    # ──────────────────────────────────────────────────────────
    for ci, company_data in enumerate(companies):
        slide    = prs.slides.add_slide(prs.slide_layouts[6])
        name     = company_data.get("name", "")
        analysis = company_data.get("analysis", {})
        articles = company_data.get("articles", [])

        add_rect(slide, 0, 0, W, Inches(1.1), AMBER)
        add_textbox(slide, Inches(0.5), Inches(0.15), Inches(8), Inches(0.7),
                    name, font_size=24, bold=True, color=WHITE)
        add_textbox(slide, Inches(11.5), Inches(0.25), Inches(1.3), Inches(0.5),
                    f"{ci + 2} / {len(companies) + 2}",
                    font_size=10, color=RGBColor(0xFF, 0xE0, 0x9A), align=PP_ALIGN.RIGHT)

        col_w       = Inches(6.2)
        col_gap     = Inches(0.3)
        row_h       = Inches(2.5)
        x_positions = [Inches(0.3), Inches(0.3) + col_w + col_gap]
        y_positions = [Inches(1.3), Inches(1.3) + row_h + Inches(0.2)]

        for i, cat in enumerate(categories):
            col = i % 2
            row = i // 2
            x   = x_positions[col]
            y   = y_positions[row] if row < len(y_positions) else y_positions[-1]

            add_rect(slide, x, y, col_w, row_h, GRAY_LITE,
                     line_color=RGBColor(0xE5, 0xE7, 0xEB))
            add_rect(slide, x, y, col_w, Inches(0.35), AMBER_LITE)
            add_textbox(slide, x + Inches(0.15), y + Inches(0.04),
                        col_w - Inches(0.3), Inches(0.28),
                        cat, font_size=10, bold=True, color=AMBER_DARK)
            content = analysis.get(cat, "해당 없음")
            add_textbox(slide, x + Inches(0.15), y + Inches(0.42),
                        col_w - Inches(0.3), row_h - Inches(0.55),
                        content, font_size=9, color=DARK, wrap=True)

        if articles:
            src_text = "참고 출처: " + " · ".join(
                a.get("title", "")[:30] for a in articles[:3]
            )
            add_textbox(slide, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.3),
                        src_text, font_size=7, color=GRAY)

    # ──────────────────────────────────────────────────────────
    # 마지막 슬라이드: 종합 시사점
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, Inches(1.1), AMBER_DARK)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(10), Inches(0.7),
                "종합 시사점 및 전략적 액션 아이템",
                font_size=22, bold=True, color=WHITE)
    add_textbox(slide, Inches(11.5), Inches(0.25), Inches(1.3), Inches(0.5),
                f"{len(companies) + 2} / {len(companies) + 2}",
                font_size=10, color=RGBColor(0xFF, 0xE0, 0x9A), align=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0.3), Inches(1.2), W - Inches(0.6), H - Inches(1.5), AMBER_LITE,
             line_color=RGBColor(0xFD, 0xD8, 0x35))
    add_textbox(slide, Inches(0.5), Inches(1.35), W - Inches(1.0), H - Inches(2.0),
                summary or "종합 시사점이 없습니다.",
                font_size=11, color=DARK, wrap=True)

    # ──────────────────────────────────────────────────────────
    # 재무 비교 차트 슬라이드 (financial_data가 있을 때만)
    # ──────────────────────────────────────────────────────────
    fin_data   = data.get("financial_data", {})
    ticker_map = data.get("ticker_map", {})

    if fin_data and ticker_map:
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE

        # 회사별 재무 데이터 매핑
        comp_fin = []
        for c in companies:
            cname = c["name"]
            ti     = ticker_map.get(cname, {})
            ticker = ti.get("ticker") if ti.get("found") else None
            fd     = fin_data.get(ticker, {}) if ticker else {}
            comp_fin.append({"name": cname, "fd": fd})

        # ── 스케일 결정 헬퍼 ──────────────────────────────────
        def _scale(vals, currency):
            """값 목록으로 (divisor, unit) 결정"""
            valid = [v for v in vals if v and v > 0]
            if not valid:
                return 1.0, ""
            mx = max(valid)
            if currency == "KRW":
                if mx >= 1e12: return 1e12, "(조원)"
                return 1e8, "(억원)"
            if mx >= 1e12: return 1e12, "(T)"
            if mx >= 1e9:  return 1e9,  "(B)"
            return 1e6, "(M)"

        # ── 차트 슬라이드 생성 헬퍼 ──────────────────────────
        def _add_chart_slide(title, series_list, no_data_names):
            """
            series_list: [(series_label, [val_or_None, ...]) per categories_list order]
            no_data_names: 데이터 없는 회사 이름 목록
            """
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 헤더
            add_rect(slide, 0, 0, W, Inches(1.1), AMBER_DARK)
            add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
                        title, font_size=20, bold=True, color=WHITE)

            # 차트 높이 — 하단 주석 유무에 따라
            chart_h = Inches(4.3) if no_data_names else Inches(4.8)

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED,
                Inches(0.5), Inches(1.25),
                W - Inches(1.0), chart_h,
                series_list,  # ChartData 객체
            )
            ch = chart_frame.chart

            # 범례: 시리즈가 2개 이상이면 표시
            try:
                ch.has_legend = len(list(ch.series)) > 1
                if ch.has_legend:
                    from pptx.enum.chart import XL_LEGEND_POSITION
                    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
                    ch.legend.include_in_layout = False
            except Exception:
                pass

            # 시리즈 색상 amber 적용
            try:
                colors = [CHART_C1, CHART_C2]
                for i, s in enumerate(ch.series):
                    s.format.fill.solid()
                    s.format.fill.fore_color.rgb = colors[i % 2]
            except Exception:
                pass

            # 데이터 없는 회사 주석
            if no_data_names:
                add_textbox(slide, Inches(0.5), H - Inches(0.6), W - Inches(1.0), Inches(0.4),
                            f"데이터 없음: {', '.join(no_data_names)}",
                            font_size=9, color=GRAY)

        # ── 차트 1: 시가총액 ──────────────────────────────────
        mc_vals     = [c["fd"].get("marketCap") for c in comp_fin]
        mc_currency = next((c["fd"].get("currency") for c in comp_fin if c["fd"].get("currency")), "")
        mc_div, mc_unit = _scale(mc_vals, mc_currency)

        has_mc  = [c for c in comp_fin if c["fd"].get("marketCap")]
        no_mc   = [c["name"] for c in comp_fin if not c["fd"].get("marketCap")]

        if has_mc:
            cd = ChartData()
            cd.categories = [c["name"] for c in has_mc]
            cd.add_series(f"시가총액 {mc_unit}",
                          tuple((c["fd"]["marketCap"] / mc_div) for c in has_mc))
            _add_chart_slide("시가총액 비교", cd, no_mc)

        # ── 차트 2: 매출 · 영업이익 ───────────────────────────
        rev_vals  = [c["fd"].get("totalRevenue") for c in comp_fin]
        oi_vals   = [c["fd"].get("operatingIncome") for c in comp_fin]
        fin_currency = next((c["fd"].get("currency") for c in comp_fin if c["fd"].get("currency")), "")
        rev_div, rev_unit = _scale(
            [v for v in rev_vals + oi_vals if v is not None],
            fin_currency,
        )

        has_rev = [c for c in comp_fin
                   if c["fd"].get("totalRevenue") or c["fd"].get("operatingIncome")]
        no_rev  = [c["name"] for c in comp_fin
                   if not c["fd"].get("totalRevenue") and not c["fd"].get("operatingIncome")]

        if has_rev:
            cd = ChartData()
            cd.categories = [c["name"] for c in has_rev]
            cd.add_series(f"매출액 {rev_unit}",
                          tuple((c["fd"].get("totalRevenue") or 0) / rev_div for c in has_rev))
            cd.add_series(f"영업이익 {rev_unit}",
                          tuple((c["fd"].get("operatingIncome") or 0) / rev_div for c in has_rev))
            _add_chart_slide("매출 · 영업이익 비교", cd, no_rev)

        # ── 차트 3: PER · PBR ─────────────────────────────────
        has_per = [c for c in comp_fin
                   if c["fd"].get("trailingPE") or c["fd"].get("priceToBook")]
        no_per  = [c["name"] for c in comp_fin
                   if not c["fd"].get("trailingPE") and not c["fd"].get("priceToBook")]

        if has_per:
            cd = ChartData()
            cd.categories = [c["name"] for c in has_per]
            cd.add_series("PER (배)",
                          tuple(c["fd"].get("trailingPE") or 0 for c in has_per))
            cd.add_series("PBR (배)",
                          tuple(c["fd"].get("priceToBook") or 0 for c in has_per))
            _add_chart_slide("PER · PBR 비교", cd, no_per)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
